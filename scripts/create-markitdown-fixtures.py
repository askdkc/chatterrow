#!/usr/bin/env python3

"""Create small real PDF, DOCX, XLSX, and PPTX fixtures for CI smoke tests."""

from html import escape
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile


ROOT = Path(__file__).resolve().parents[1]
FIXTURE_DIR = ROOT / "tests" / "Fixtures" / "markitdown"


def pdf_document(text: str) -> bytes:
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream = f"BT /F1 18 Tf 72 720 Td ({escaped}) Tj ET".encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    document = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]

    for number, content in enumerate(objects, start=1):
        offsets.append(len(document))
        document.extend(f"{number} 0 obj\n".encode("ascii"))
        document.extend(content)
        document.extend(b"\nendobj\n")

    xref = len(document)
    document.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    document.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        document.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    document.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref}\n%%EOF\n"
        ).encode("ascii")
    )

    return bytes(document)


def create_docx() -> None:
    document_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    <w:p><w:pPr><w:pStyle w:val="Title"/></w:pPr><w:r><w:t>DOCX conversion fixture</w:t></w:r></w:p>
    <w:p><w:r><w:t>{escape("DOCX conversion fixture text.")}</w:t></w:r></w:p>
    <w:sectPr>
      <w:pgSz w:w="12240" w:h="15840"/>
      <w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440"/>
    </w:sectPr>
  </w:body>
</w:document>
'''
    content_types = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>
'''
    package_relationships = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>
'''

    with ZipFile(FIXTURE_DIR / "document.docx", "w", ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", package_relationships)
        archive.writestr("word/document.xml", document_xml)


def create_xlsx() -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Fixture"
    worksheet.append(["XLSX conversion fixture", "XLSX conversion fixture text."])
    workbook.save(FIXTURE_DIR / "document.xlsx")


def create_pptx() -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "PPTX conversion fixture"
    slide.placeholders[1].text = "PPTX conversion fixture text."
    presentation.save(FIXTURE_DIR / "document.pptx")


def main() -> None:
    FIXTURE_DIR.mkdir(parents=True, exist_ok=True)
    (FIXTURE_DIR / "document.pdf").write_bytes(pdf_document("PDF conversion fixture text."))
    create_docx()
    create_xlsx()
    create_pptx()


if __name__ == "__main__":
    main()
